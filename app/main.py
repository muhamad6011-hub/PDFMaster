from fastapi import FastAPI, Request, Form
from fastapi.responses import HTMLResponse, RedirectResponse
from fastapi.templating import Jinja2Templates
from starlette.middleware.sessions import SessionMiddleware

from sqlalchemy import create_engine, Column, Integer, String, or_
from sqlalchemy.orm import sessionmaker, declarative_base

from passlib.context import CryptContext
from dotenv import load_dotenv
from docx2pdf import convert
from PyPDF2 import PdfReader, PdfWriter
from pdf2docx import Converter
from PIL import Image
from pdf2image import convert_from_path

import uuid
import pandas as pd
from reportlab.pdfgen import canvas
from pptx import Presentation


import os


# ================= LOAD ENV =================

load_dotenv()
SECRET_KEY = os.getenv("SECRET_KEY", "super-secret-key")

# ================= APP =================

app = FastAPI()

from fastapi.staticfiles import StaticFiles

app.mount("/static", StaticFiles(directory="app/static"), name="static")

app.add_middleware(
    SessionMiddleware,
    secret_key=SECRET_KEY
)

# ================= TEMPLATES =================

templates = Jinja2Templates(directory="app/templates")

# ================= DATABASE =================

DATABASE_URL = "sqlite:///./database.db"

engine = create_engine(
    DATABASE_URL,
    connect_args={"check_same_thread": False}
)

SessionLocal = sessionmaker(bind=engine)
Base = declarative_base()

pwd_context = CryptContext(
    schemes=["pbkdf2_sha256"],
    deprecated="auto"
)

class User(Base):
    __tablename__ = "users"

    id = Column(Integer, primary_key=True, index=True)
    username = Column(String, unique=True)
    email = Column(String, unique=True)
    password = Column(String)

class BlogPost(Base):
    __tablename__ = "blog_posts"

    id = Column(Integer, primary_key=True, index=True)
    title = Column(String)
    slug = Column(String, unique=True, index=True)
    content = Column(String)

# create tables
Base.metadata.create_all(bind=engine)

# ================= HOME =================

@app.get("/", response_class=HTMLResponse)
async def index(request: Request):
    return templates.TemplateResponse("index.html", {"request": request})

# ================= TOOLS =================

@app.get("/tools/{tool_name}", response_class=HTMLResponse)
async def open_tool(request: Request, tool_name: str):

    if not request.session.get("user"):
        return RedirectResponse("/login", status_code=302)

    template_name = tool_name.replace("-", "_")

    return templates.TemplateResponse(
        f"tools/{template_name}.html",
        {"request": request}
    )

from fastapi import UploadFile, File
from fastapi.responses import FileResponse
import shutil


UPLOAD_FOLDER = "app/uploads"
OUTPUT_FOLDER = "app/outputs"

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

# ================= COMPRESS PDF (REAL TARGET) =================

import subprocess

@app.post("/tools/compress-pdf")
async def compress_pdf(
    request: Request,
    file: UploadFile = File(...),
    target_size: int = Form(...)
):

    if not request.session.get("user"):
        return RedirectResponse("/login", status_code=302)

    input_path = os.path.join(UPLOAD_FOLDER, file.filename)
    output_path = os.path.join(
        OUTPUT_FOLDER,
        f"compressed_{file.filename}"
    )

    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    # Tentukan DPI berdasarkan target
    if target_size <= 100:
        dpi = 50
    elif target_size <= 200:
        dpi = 72
    elif target_size <= 300:
        dpi = 100
    else:
        dpi = 150

    command = [
        "gswin64c",
        "-sDEVICE=pdfwrite",
        "-dCompatibilityLevel=1.4",
        "-dNOPAUSE",
        "-dQUIET",
        "-dBATCH",
        "-dDownsampleColorImages=true",
        "-dDownsampleGrayImages=true",
        "-dDownsampleMonoImages=true",
        f"-dColorImageResolution={dpi}",
        f"-dGrayImageResolution={dpi}",
        f"-dMonoImageResolution={dpi}",
        f"-sOutputFile={output_path}",
        input_path,
    ]

    subprocess.run(command)

    return FileResponse(
        output_path,
        filename=f"compressed_{file.filename}",
        media_type="application/pdf"
    )
# ================= WORD TO PDF =================

@app.post("/tools/word-to-pdf")
async def word_to_pdf(
    request: Request,
    file: UploadFile = File(...)
):

    if not request.session.get("user"):
        return RedirectResponse("/login", status_code=302)

    input_path = os.path.join(UPLOAD_FOLDER, file.filename)
    output_path = os.path.join(
        OUTPUT_FOLDER,
        file.filename.replace(".docx", ".pdf")
    )

    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    convert(input_path, output_path)

    return FileResponse(
        output_path,
        filename="converted.pdf",
        media_type="application/pdf"
    )

# ================= JPG TO PDF =================

@app.post("/tools/jpg-to-pdf")
async def jpg_to_pdf(request: Request, file: UploadFile = File(...)):

    if not request.session.get("user"):
        return RedirectResponse("/login", status_code=302)

    input_path = os.path.join(UPLOAD_FOLDER, file.filename)
    output_path = os.path.join(OUTPUT_FOLDER, "result.pdf")

    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    image = Image.open(input_path)
    image.convert("RGB").save(output_path)

    return FileResponse(output_path, filename="result.pdf")

# ================= MERGE PDF =================

@app.post("/tools/merge-pdf")
async def merge_pdf(
    request: Request,
    files: list[UploadFile] = File(...)
):

    if not request.session.get("user"):
        return RedirectResponse("/login", status_code=302)

    if len(files) < 2:
        return RedirectResponse("/tools/merge-pdf", status_code=302)

    writer = PdfWriter()

    for file in files:
        input_path = os.path.join(UPLOAD_FOLDER, file.filename)

        # Simpan file upload
        with open(input_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        reader = PdfReader(input_path)

        # Tambahkan semua halaman ke writer
        for page in reader.pages:
            writer.add_page(page)

    output_path = os.path.join(OUTPUT_FOLDER, "merged.pdf")

    with open(output_path, "wb") as f:
        writer.write(f)

    return FileResponse(
        output_path,
        filename="merged.pdf",
        media_type="application/pdf"
    )

# ================= PDF → JPG =================

@app.post("/tools/pdf-to-jpg")
async def pdf_to_jpg(request: Request, file: UploadFile = File(...)):

    if not request.session.get("user"):
        return RedirectResponse("/login", status_code=302)

    input_path = os.path.join(UPLOAD_FOLDER, file.filename)

    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    images = convert_from_path(input_path)

    output_path = os.path.join(OUTPUT_FOLDER, "page.jpg")
    images[0].save(output_path, "JPEG")

    return FileResponse(output_path, filename="page.jpg")

# ================= PDF → WORD =================

@app.post("/tools/pdf-to-word")
async def pdf_to_word(request: Request, file: UploadFile = File(...)):

    if not request.session.get("user"):
        return RedirectResponse("/login", status_code=302)

    input_path = os.path.join(UPLOAD_FOLDER, file.filename)
    output_path = os.path.join(
        OUTPUT_FOLDER,
        file.filename.replace(".pdf", ".docx")
    )

    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    cv = Converter(input_path)
    cv.convert(output_path)
    cv.close()

    return FileResponse(output_path, filename="result.docx")

# ================= SPLIT PDF =================

@app.post("/tools/split-pdf")
async def split_pdf(request: Request, file: UploadFile = File(...)):

    if not request.session.get("user"):
        return RedirectResponse("/login", status_code=302)

    input_path = os.path.join(UPLOAD_FOLDER, file.filename)

    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    reader = PdfReader(input_path)

    writer = PdfWriter()

    for page in reader.pages:
        writer.add_page(page)

    output_path = os.path.join(OUTPUT_FOLDER, "split.pdf")

    with open(output_path, "wb") as f:
        writer.write(f)

    return FileResponse(output_path, filename="split.pdf")

# ================= COMPRESS IMAGE =================

@app.post("/tools/compress-image")
async def compress_image(
    request: Request,
    file: UploadFile = File(...),
    quality: int = Form(70)
):

    filename = f"{uuid.uuid4()}_{file.filename}"

    input_path = os.path.join(UPLOAD_FOLDER, filename)
    output_path = os.path.join(OUTPUT_FOLDER, f"compressed_{filename}")

    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    image = Image.open(input_path)

    image.save(
        output_path,
        optimize=True,
        quality=quality
    )

    return FileResponse(
        output_path,
        filename=f"compressed_{file.filename}"
    )
# ================= PNG TO PDF =================

@app.post("/tools/png-to-pdf")
async def png_to_pdf(
    request: Request,
    file: UploadFile = File(...)
):

    if not request.session.get("user"):
        return RedirectResponse("/login", status_code=302)

    filename = f"{uuid.uuid4()}_{file.filename}"

    input_path = os.path.join(UPLOAD_FOLDER, filename)
    output_path = os.path.join(OUTPUT_FOLDER, f"{filename}.pdf")

    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    image = Image.open(input_path)

    image.convert("RGB").save(output_path)

    return FileResponse(
        output_path,
        filename="converted.pdf"
    )
# ================= EXCEL TO PDF =================

@app.post("/tools/excel-to-pdf")
async def excel_to_pdf(
    request: Request,
    file: UploadFile = File(...)
):

    if not request.session.get("user"):
        return RedirectResponse("/login", status_code=302)

    filename = f"{uuid.uuid4()}_{file.filename}"

    input_path = os.path.join(UPLOAD_FOLDER, filename)
    output_path = os.path.join(OUTPUT_FOLDER, f"{filename}.pdf")

    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    df = pd.read_excel(input_path)

    c = canvas.Canvas(output_path)

    y = 800

    for row in df.values:
        text = " | ".join(str(x) for x in row)
        c.drawString(50, y, text)
        y -= 20

    c.save()

    return FileResponse(
        output_path,
        filename="excel.pdf"
    )
# ================= PPT TO PDF =================

@app.post("/tools/ppt-to-pdf")
async def ppt_to_pdf(
    request: Request,
    file: UploadFile = File(...)
):

    if not request.session.get("user"):
        return RedirectResponse("/login", status_code=302)

    filename = f"{uuid.uuid4()}_{file.filename}"

    input_path = os.path.join(UPLOAD_FOLDER, filename)
    output_path = os.path.join(OUTPUT_FOLDER, f"{filename}.pdf")

    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    prs = Presentation(input_path)

    c = canvas.Canvas(output_path)

    y = 800

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                c.drawString(50, y, shape.text)
                y -= 20

        c.showPage()

    c.save()

    return FileResponse(
        output_path,
        filename="presentation.pdf"
    )


# ================= REGISTER =================

@app.get("/register", response_class=HTMLResponse)
async def register_page(request: Request):
    return templates.TemplateResponse("register.html", {"request": request})

@app.post("/register")
async def register(
    request: Request,
    username: str = Form(...),
    email: str = Form(...),
    password: str = Form(...)
):

    db = SessionLocal()

    existing_user = db.query(User).filter(
        or_(User.username == username, User.email == email)
    ).first()

    if existing_user:
        db.close()
        return templates.TemplateResponse(
            "register.html",
            {"request": request, "error": "Username or email already registered"}
        )

    hashed_password = pwd_context.hash(password)

    new_user = User(
        username=username,
        email=email,
        password=hashed_password
    )

    db.add(new_user)
    db.commit()
    db.close()

    return RedirectResponse("/login", status_code=302)

# ================= LOGIN =================

@app.get("/login", response_class=HTMLResponse)
async def login_page(request: Request):
    return templates.TemplateResponse("login.html", {"request": request})

@app.post("/login")
async def login(
    request: Request,
    username: str = Form(...),
    password: str = Form(...)
):

    db = SessionLocal()
    user = db.query(User).filter(User.username == username).first()

    if not user or not pwd_context.verify(password, user.password):
        db.close()
        return templates.TemplateResponse(
            "login.html",
            {"request": request, "error": "Invalid username or password"}
        )

    request.session["user"] = user.username
    db.close()

    return RedirectResponse("/dashboard", status_code=302)

# ================= DASHBOARD =================

@app.get("/dashboard", response_class=HTMLResponse)
async def dashboard(request: Request):

    if not request.session.get("user"):
        return RedirectResponse("/login", status_code=302)

    return templates.TemplateResponse(
        "dashboard.html",
        {"request": request}
    )

# ================= ADMIN CREATE POST PAGE =================

@app.get("/admin/create-post", response_class=HTMLResponse)
def create_post_page(request: Request):

    if request.session.get("user") != "admin":
        return RedirectResponse("/", status_code=302)

    return templates.TemplateResponse(
        "admin/create_post.html",
        {"request": request}
    )
# ================= ADMIN SAVE POST =================

@app.post("/admin/create-post")
def create_post(
    request: Request,
    title: str = Form(...),
    slug: str = Form(...),
    content: str = Form(...)
):

    if request.session.get("user") != "admin":
        return RedirectResponse("/", status_code=302)

    db = SessionLocal()

    post = BlogPost(
        title=title,
        slug=slug,
        content=content
    )

    db.add(post)
    db.commit()
    db.close()

    return RedirectResponse("/admin/posts", status_code=302)

# ================= ADMIN LIST POSTS =================
@app.get("/admin/posts", response_class=HTMLResponse)
def admin_posts(request: Request):

    if request.session.get("user") != "admin":
        return RedirectResponse("/", status_code=302)

    db = SessionLocal()

    posts = db.query(BlogPost)\
        .order_by(BlogPost.id.desc())\
        .all()

    db.close()

    return templates.TemplateResponse(
        "admin/posts.html",
        {
            "request": request,
            "posts": posts
        }
    )

# ================= ADMIN DELETE POST =================
@app.get("/admin/delete-post/{post_id}")
def delete_post(post_id: int, request: Request):

    if request.session.get("user") != "admin":
        return RedirectResponse("/", status_code=302)

    db = SessionLocal()

    post = db.query(BlogPost)\
        .filter(BlogPost.id == post_id)\
        .first()

    if post:
        db.delete(post)
        db.commit()

    db.close()

    return RedirectResponse("/admin/posts", status_code=302)
# ================= BLOG LIST =================
@app.get("/blog", response_class=HTMLResponse)
def blog(request: Request):

    db = SessionLocal()

    posts = db.query(BlogPost)\
        .order_by(BlogPost.id.desc())\
        .all()

    db.close()

    return templates.TemplateResponse(
        "blog.html",
        {
            "request": request,
            "posts": posts
        }
    )
# ================= BLOG POST =================
@app.get("/blog/{slug}", response_class=HTMLResponse)
def blog_post(slug: str, request: Request):

    db = SessionLocal()

    post = db.query(BlogPost)\
        .filter(BlogPost.slug == slug)\
        .first()

    db.close()

    if not post:
        return RedirectResponse("/blog", status_code=302)

    return templates.TemplateResponse(
        "blog_post.html",
        {
            "request": request,
            "post": post
        }
    )
# ================= STATIC PAGES =================
@app.get("/about")
def about(request: Request):
    return templates.TemplateResponse("about.html", {"request": request})


@app.get("/contact")
def contact(request: Request):
    return templates.TemplateResponse("contact.html", {"request": request})


@app.get("/privacy")
def privacy(request: Request):
    return templates.TemplateResponse("privacy.html", {"request": request})


@app.get("/terms")
def terms(request: Request):
    return templates.TemplateResponse("terms.html", {"request": request})


@app.get("/disclaimer")
def disclaimer(request: Request):
    return templates.TemplateResponse("disclaimer.html", {"request": request})

# ================= LOGOUT =================

@app.get("/logout")
async def logout(request: Request):
    request.session.clear()
    return RedirectResponse("/", status_code=302)

# ================= SITEMAP =================

@app.get("/sitemap.xml")
def sitemap():

    db = SessionLocal()
    posts = db.query(BlogPost).all()
    db.close()

    urls = ""

    for post in posts:
        urls += f"""
        <url>
            <loc>https://pdfmaster.com/blog/{post.slug}</loc>
        </url>
        """

    xml = f"""
    <urlset xmlns="http://www.sitemaps.org/schemas/sitemap/0.9">

        <url>
            <loc>https://pdfmaster.com/</loc>
        </url>

        <url>
            <loc>https://pdfmaster.com/blog</loc>
        </url>

        {urls}

    </urlset>
    """

    return HTMLResponse(content=xml, media_type="application/xml")

@app.get("/debug-blog")
def debug_blog():

    db = SessionLocal()
    posts = db.query(BlogPost).all()
    db.close()

    return {"total_posts": len(posts)}

@app.get("/debug-blog")
def debug_blog():

    db = SessionLocal()
    posts = db.query(BlogPost).all()
    db.close()

    return {"total_posts": len(posts)}
